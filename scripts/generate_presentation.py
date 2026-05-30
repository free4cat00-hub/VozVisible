#!/usr/bin/env python3
"""Generador simple de presentación .pptx para VozVisible.

Crea un PPTX con la estructura acordada y añade notas del ponente.
Requiere: python-pptx, pillow

Uso:
    pip install python-pptx pillow
    python scripts/generate_presentation.py

El archivo resultante se guarda en `assets/presentation/VozVisible_presentation.pptx`.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "assets" / "presentation"
OUT_DIR.mkdir(parents=True, exist_ok=True)

VIDEOS = [
    ROOT / "assets" / "output" / "el_tren_test_regen.mp4",
    ROOT / "spoken_to_signed" / "videos_lse" / "tren.mp4",
    ROOT / "spoken_to_signed" / "videos_lse" / "salir.mp4",
]


def make_thumbnail(video_path: Path, out_path: Path, size=(960, 540)):
    # Simple placeholder thumbnail with filename text (no ffmpeg dependency)
    img = Image.new("RGB", size, color=(30, 30, 30))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans-Bold.ttf", 32)
    except Exception:
        font = ImageFont.load_default()
    text = video_path.name
    w, h = draw.textsize(text, font=font)
    draw.text(((size[0]-w)/2, (size[1]-h)/2), text, font=font, fill=(235, 235, 235))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Duración: 5 minutos presentación + 2 preguntas. Intro rápida y pitch." 


def add_bullet_slide(prs: Presentation, title: str, bullets, notes_text: str = ""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def add_demo_slide(prs: Presentation, video_thumbs):
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    title_box = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.5))
    tf = title_box.text_frame
    tf.text = "Demo corta"
    tf.paragraphs[0].font.size = Pt(28)

    thumb_w = Inches(3)
    thumb_h = Inches(1.9)
    x = left
    y = Inches(1)
    for t in video_thumbs:
        if t.exists():
            slide.shapes.add_picture(str(t), x, y, width=thumb_w, height=thumb_h)
        else:
            tb = slide.shapes.add_textbox(x, y, thumb_w, thumb_h)
            tb.text = f"Vídeo: {t.name} (no encontrado)"
        x += thumb_w + Inches(0.2)
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Mostrar 30-60s: flujo de texto → gloss → pose → render. "
        "Archivos de vídeo locales enlazados en la carpeta `assets/videos` y `assets/output`."
    )


def build_presentation():
    prs = Presentation()

    add_title_slide(prs, "VozVisible — Traducción hablada a LSE", "Demo y propuesta de proyecto")

    add_bullet_slide(
        prs,
        "Problema / Oportunidad",
        [
            "Información accesible para personas sordas y con dificultades auditivas es limitada.",
            "Escasez de soluciones automáticas de LSE en entornos reales (transporte, educación).",
            "Gran mercado y potencial de integración en servicios públicos y privados.",
        ],
        "Breve introducción del problema, conectar con mercado y oportunidad de venta."
    )

    add_demo_slide(prs, [OUT_DIR / (v.name + ".thumb.png") for v in VIDEOS])

    add_bullet_slide(
        prs,
        "Cómo funciona (pipeline)",
        [
            "Audio → Speech-to-text (ASR)",
            "Texto → Gloss / tokenización LSE (reglas + LLM)",
            "Gloss → Pose (lexicón / modelos)",
            "Pose → Render de vídeo (pose-to-video)",
        ],
        "Mencionar componentes: Spacy, transformers, pose-to-video, ffmpeg para render." 
    )

    add_bullet_slide(
        prs,
        "Tracción y modelo",
        [
            "Dataset inicial: lexicones y vídeos de demostración (assets).",
            "Prueba local completa; demo listo para mostrar en 5 minutos.",
            "Modelo de negocio: licencia a instituciones, SaaS para servicios, integraciones B2B.",
        ],
        "Incluir métricas si hay (usuarios, tests, accuracy) y próximos pasos para validación."
    )

    add_bullet_slide(
        prs,
        "Roadmap & Equipo",
        [
            "MVP estable en local → despliegue en servidor → demo pública.",
            "Mejoras: robustez en streaming, fallback transcode, UI/UX, datasets ampliados.",
            "Equipo: devs (tú), integrador ASR, diseñador vídeo, búsqueda de mentores/inversión.",
        ],
        "Plazos estimados: 3 meses para MVP público; 6-9 meses para producto comercializable."
    )

    add_bullet_slide(
        prs,
        "Cierre / CTA",
        [
            "Demo en vivo disponible; paquete técnico y pruebas bajo demanda.",
            "Buscamos feedback, pilotos y posibles inversores/partners.",
        ],
        "Contacto y llamada a la acción: demo extendida, solicitar reunión inversores."
    )

    # Prepare thumbnails for demo videos
    thumbs = []
    for v in VIDEOS:
        thumb = OUT_DIR / (v.name + ".thumb.png")
        make_thumbnail(v, thumb)
        thumbs.append(thumb)

    # overwrite demo slide pictures with generated thumbs
    # (we already used the thumb filenames when creating the demo slide)

    out_path = OUT_DIR / "VozVisible_presentation.pptx"
    prs.save(str(out_path))
    print(f"Presentación generada: {out_path}")


if __name__ == "__main__":
    build_presentation()
